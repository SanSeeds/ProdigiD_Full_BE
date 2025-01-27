import json
import logging
import os
import random
from PIL import Image
from django.conf import settings
from groq import Groq
import requests
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from better_profanity import profanity
from django.http import JsonResponse, HttpResponse
from docx import Document as DocxDocument
import fitz  # PyMuPDF
import openpyxl
import xlrd
from langchain_core.output_parsers import StrOutputParser
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnablePassthrough
from langchain_google_genai import GoogleGenerativeAIEmbeddings
# from langchain_astradb import AstraDBVectorStore
from langchain.indexes import VectorstoreIndexCreator
from langchain_groq import ChatGroq
import base64

# from prodigidusk.settings import VECTOR_STORE

logger = logging.getLogger(__name__)

# API Keys
GROQ_SECRET_ACCESS_KEY = settings.GROQ_SECRET_ACCESS_KEY
BHASHINI_API_KEY = settings.BHASHINI_API_KEY
BHASHINI_USER_ID = settings.BHASHINI_USER_ID
VECTOR_STORE = settings.FAISS_VECTOR_STORE


# Function to check for inappropriate language
def contains_inappropriate_language(text: str) -> bool:
    return profanity.contains_profanity(text)

# Function to sanitize input containing inappropriate words
def sanitize_input(input_str):
    return profanity.censor(input_str)



MAX_SENTENCES_PER_CHUNK = 30  # Number of sentences to process in a single chunk
MAX_WORKERS = 50  # Number of threads for concurrent processing
RETRY_LIMIT = 1000  # Maximum retries for translation API
PEXELS_API_KEY = "WjL9jZR3XAwZOE1zt7laoApHMSKfYwuDmUQVbKEaYRglJg18FcfrtiG8"
PEXELS_BASE_URL = "https://api.pexels.com/v1/search"
 
# Function to translate a single sentence
def bhashini_translate_formatted(sentence, to_code, from_code="English"):
    lang_dict = {
        "English": "en",
        "Hindi": "hi",
        "Tamil": "ta",
        "Telugu": "te",
        "Marathi": "mr",
        "Kannada": "kn",
        "Bengali": "bn",
        "Odia": "or",
        "Assamese": "as",
        "Punjabi": "pa",
        "Malayalam": "ml",
        "Gujarati": "gu",
        "Urdu": "ur",
        "Sanskrit": "sa",
        "Nepali": "ne",
        "Bodo": "brx",
        "Maithili": "mai",
        "Sindhi": "sd",
        "Kashmiri": "ks",
        "Konkani": "kok",
        "Dogri": "doi",
        "Goan Konkani": "gom",
        "Santali": "sat",
    }

    from_code = lang_dict[from_code]
    to_code = lang_dict[to_code]

    url = 'https://meity-auth.ulcacontrib.org/ulca/apis/v0/model/getModelsPipeline'
    headers = {
        "Content-Type": "application/json",
        "userID": BHASHINI_USER_ID,
        "ulcaApiKey": BHASHINI_API_KEY
    }
    payload = {
        "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}}}],
        "pipelineRequestConfig": {"pipelineId": "64392f96daac500b55c543cd"}
    }

    try:
        response = requests.post(url, json=payload, headers=headers, timeout=10)  # 10s timeout
        response.raise_for_status()  # Raise exception for non-200 status codes
        response_data = response.json()

        # Prepare translation request
        service_id = response_data["pipelineResponseConfig"][0]["config"][0]["serviceId"]
        callback_url = response_data["pipelineInferenceAPIEndPoint"]["callbackUrl"]
        headers2 = {
            "Content-Type": "application/json",
            response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["name"]: response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["value"]
        }
        compute_payload = {
            "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}, "serviceId": service_id}}],
            "inputData": {"input": [{"source": sentence}], "audio": [{"audioContent": None}]}
        }

        compute_response = requests.post(callback_url, json=compute_payload, headers=headers2, timeout=10)  # 10s timeout
        compute_response.raise_for_status()
        compute_response_data = compute_response.json()
        return compute_response_data["pipelineResponse"][0]["output"][0]["target"]
    except requests.RequestException as e:
        logger.error(f"Translation API error: {str(e)}")
        return None

# Function to handle retries for translation
def translate_with_retry(sentence, to_code, retries=RETRY_LIMIT):
    for attempt in range(retries):
        result = bhashini_translate_formatted(sentence, to_code)
        if result:
            return result
    raise ValueError(f"Failed to translate sentence after {retries} attempts: {sentence}")

def generate_email(purpose='Request Information', num_words=100, subject=None, rephrase=False, to=None, tone='Formal', keywords=None, contextual_background=None, call_to_action=None, additional_details=None, priority_level='Low', closing_remarks=None):
    # Ensure all fields are checked for inappropriate language
    fields_to_check = [purpose, subject, keywords, contextual_background, call_to_action, additional_details, closing_remarks]
    if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
        return "Error: Input contains inappropriate language."
    
    # Construct the prompt for email generation
    prompt = f"Compose a thoughtful and engaging email of up to {num_words} words with the subject: {subject}, addressed to {to}, in a {tone} tone.\n\nIncorporate:\n"

    if purpose:
        prompt += f"- Purpose: {purpose}\n"
    if contextual_background:
        prompt += f"- Context: {contextual_background}\n"
    if call_to_action:
        prompt += f"- Call to action: {call_to_action}\n"
    if additional_details:
        prompt += f"- Additional details: {additional_details}\n"
    if priority_level:
        prompt += f"- Priority: {priority_level}\n"
    if closing_remarks:
        prompt += f"- Closing remarks: {closing_remarks}\n"

    if rephrase:
        prompt += "\nRephrase the subject line for clarity and impact if needed."

    prompt += "\nKeep the email empathetic,conversational , audience-focused, and actionable without adding extra commentary."

    # Generate content using Groq API
    client = Groq(api_key= "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")

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.3-70b-versatile",
    )

    return chat_completion.choices[0].message.content 

def generate_bus_pro(business_intro, proposal_objective, num_words, scope_of_work, project_phases, expected_outcomes, tech_innovations, target_audience, budget_info, timeline, benefits, closing_remarks):
    # Collect all fields to check for inappropriate language
    fields_to_check = [business_intro, proposal_objective, scope_of_work, project_phases, expected_outcomes, tech_innovations, target_audience, budget_info, timeline, benefits, closing_remarks]
    
    # Check if any field contains inappropriate language
    if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
        return "Error: Input contains inappropriate language."
    
    # Sanitize input (if needed)
    sanitized_fields = [sanitize_input(str(field)) if field else '' for field in fields_to_check]

    # Reconstruct the prompt with sanitized input
    
    prompt = f"Write a professional business proposal with a well-structured format. The proposal must strictly contain at least {num_words} words and may exceed slightly to ensure coherence and completeness. Follow this structure:\n\n"

    # Add inputs into the prompt
    if sanitized_fields[0]:
        prompt += f"**Introduction**: {sanitized_fields[0]}.\n"
    if sanitized_fields[1]:
        prompt += f"**Objective**: The purpose of this proposal is {sanitized_fields[1]}.\n"
    if sanitized_fields[2]:
        prompt += f"**Scope of Work**: {sanitized_fields[2]}.\n"
    if sanitized_fields[3]:
        prompt += f"**Project Phases**: The project will proceed as follows: {sanitized_fields[3]}.\n"
    if sanitized_fields[4]:
        prompt += f"**Expected Outcomes**: The anticipated results include {sanitized_fields[4]}.\n"
    if sanitized_fields[5]:
        prompt += f"**Technologies and Innovations**: Highlighting these key technologies and innovative approaches: {sanitized_fields[5]}.\n"
    if sanitized_fields[6]:
        prompt += f"**Target Audience**: This proposal is tailored for {sanitized_fields[6]}.\n"
    if sanitized_fields[7]:
        prompt += f"**Budget**: Provide an estimate or breakdown: {sanitized_fields[7]}.\n"
    if sanitized_fields[8]:
        prompt += f"**Timeline**: The expected duration and milestones are: {sanitized_fields[8]}.\n"
    if sanitized_fields[9]:
        prompt += f"**Benefits to Recipient**: These benefits are offered: {sanitized_fields[9]}.\n"
    if sanitized_fields[10]:
        prompt += f"**Closing Remarks**: {sanitized_fields[10]}.\n"

    # Add final instructions for tone and presentation
    prompt += (
        "\nEnsure the proposal is professional, concise, and action-oriented. "
        "Use clear headings and subheadings for easy readability. Avoid generic language, "
        "and make the proposal feel tailored to the recipient's needs. "
        "Focus on measurable outcomes, value propositions, and collaboration opportunities."
    )

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile",

    )

    return chat_completion.choices[0].message.content


def generate_offer_letter(company_details, candidate_name, position_title, department, status, location,
                          start_date, compensation_benefits, work_hours, terms, acceptance_deadline,
                          contact_info, documents_needed, closing_remarks):
    # Collect all fields to check for inappropriate language
    fields_to_check = [
        company_details, candidate_name, position_title, department, status, location,
        start_date, compensation_benefits, work_hours, terms, acceptance_deadline,
        contact_info, documents_needed, closing_remarks
    ]
    
    # Check if any field contains inappropriate language
    inappropriate_key = None
    inappropriate_value = None
    for value in fields_to_check:
        if value and contains_inappropriate_language(value):
            inappropriate_key = value
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in the input: {inappropriate_value}"

    # Sanitize input fields
    sanitized_fields = [sanitize_input(str(field)) if field else '' for field in fields_to_check]

    # Build the prompt with sanitized inputs
    prompt = "Generate an offer letter given the following inputs: "
    if sanitized_fields[0]:
        prompt += f"\nOur business details are {sanitized_fields[0]}, "
    if sanitized_fields[1]:
        prompt += f"\nCandidate name is {sanitized_fields[1]}, "
    if sanitized_fields[2]:
        prompt += f"\nfor the position of {sanitized_fields[2]}, "
    if sanitized_fields[3]:
        prompt += f"\nin the department: {sanitized_fields[3]}, "
    if sanitized_fields[4]:
        prompt += f"\nas a {sanitized_fields[4]} employee. "
    if sanitized_fields[5]:
        prompt += f"\nExpected to work from: {sanitized_fields[5]}. "
    if sanitized_fields[6]:
        prompt += f"\nCandidate to join on {sanitized_fields[6]}. "
    if sanitized_fields[7]:
        prompt += f"\nCandidate will receive the following compensation and benefits: {sanitized_fields[7]}. "
    if sanitized_fields[8]:
        prompt += f"\nExpected working hours: {sanitized_fields[8]}. "
    if sanitized_fields[9]:
        prompt += f"\nFollowing are the terms of the offer: {sanitized_fields[9]}. "
    if sanitized_fields[10]:
        prompt += f"\nThe last day for accepting the offer is: {sanitized_fields[10]}. "
    if sanitized_fields[11]:
        prompt += f"\nIn case of any queries contact: {sanitized_fields[11]}. "
    if sanitized_fields[12]:
        prompt += f"\nFollowing documents to be produced on the day of joining: {sanitized_fields[12]}. "
    if sanitized_fields[13]:
        prompt += f"\nIncorporate the following closing remarks in the offer letter: {sanitized_fields[13]}. "
    
    prompt += "\nThe offer letter should have a professional tone, clearly stating the terms of employment, the company's expectations, and the candidate's role."
    prompt += "\nIt should express excitement and appreciation for the candidate’s potential contribution to the company."
    prompt += "\nThe letter should include any important steps or actions the candidate needs to take to accept the offer and add acceptance statement at the bottom "
    prompt += "\nEnsure that the letter is empathetic, welcoming, and professional."
    prompt += "\nThe letter should make the candidate feel valued and appreciated for their potential to contribute to the company's success."
    
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.3-70b-versatile",

    )

    return chat_completion.choices[0].message.content

import langid

indian_languages = {
    "English": "en",
    "Hindi": "hi",
    "Tamil": "ta",
    "Telugu": "te",
    "Marathi": "mr",
    "Kannada": "kn",
    "Bengali": "bn",
    "Odia": "or",
    "Assamese": "as",
    "Punjabi": "pa",
    "Malayalam": "ml",
    "Gujarati": "gu",
    "Urdu": "ur",
    "Sanskrit": "sa",
    "Nepali": "ne",
    "Bodo": "brx",
    "Maithili": "mai",
    "Sindhi": "sd",
    "Kashmiri": "ks",
    "Konkani": "kok",
    "Dogri": "doi",
    "Goan Konkani": "gom",
    "Santali": "sat",
}

def generate_summary(document_context, main_subject, summary_purpose, length_detail, important_elements, audience, tone, format, additional_instructions, document_file=None, text=None):
    # Check for document content source
    if document_file:
        try:
            # Extract the document content from the uploaded file
            document_content = extract_document_content(document_file)
            
            print(document_content)
        except Exception as e:
            return f"Error: Could not extract content from the uploaded document. Details: {str(e)}"

        if not document_content:
            return "Error: Document is empty or could not be read!"
    elif text:
        document_content = text
        print(f"Using provided text for summarization: {document_content[:100]}...")  # Log only the first 100 characters
    else:
        return "Error: No valid document or text content provided for summarization."

    # Detect the language of the document content
    detected_lang, _ = langid.classify(document_content)
    
    # Find the language name from the code
    detected_language_name = next((name for name, code in indian_languages.items() if code == detected_lang), "Unknown")

    if detected_language_name != "English":
        # Translate the document content to English
        translation_result = bhashini_translate(document_content, to_code="English", from_code=detected_language_name)
        if translation_result["status_code"] != 200:
            return f"Error in translation: {translation_result['message']}"
        document_content = translation_result["translated_content"]
        print(f"Translated Content: {document_content[:100]}...")  # Log only the first 100 characters

    # Count the number of words in the document
    word_count = len(document_content.split())
    if word_count > 20000:
        return "Error: Content too large. Please provide content with fewer than 20,000 words."

    # Collect all fields to check for inappropriate language
    inputs = {
        "Document Context": document_context,
        "Main Subject": main_subject,
        "Summary Purpose": summary_purpose,
        "Length Detail": length_detail,
        "Important Elements": important_elements,
        "Audience": audience,
        "Tone": tone,
        "Format": format,
        "Additional Instructions": additional_instructions
    }

    # Check if any input parameter contains inappropriate words
    inappropriate_key = None
    inappropriate_value = None
    for key, value in inputs.items():
        if value and contains_inappropriate_language(value):
            inappropriate_key = key
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in {inappropriate_key}: {inappropriate_value}"

    # Build the prompt with the inputs
    prompt = f"Please summarize the following document content based on the provided instructions:\n\n"
    prompt += f"Document Content: {document_content}\n\n"
    prompt += "Summary Instructions:\n"
    for key, value in inputs.items():
        prompt += f"- {key}: {value}\n"

    if length_detail == 'Brief Summary':
        prompt += "Keep it concise, around 150 words."
    elif length_detail == 'Standard Summary':
        prompt += "Provide a detailed summary, around 300 words."
    else:  # In-Depth Summary
        prompt += "Offer a comprehensive summary, around 450 words."

    prompt += (
        "\n\nPlease ensure the summary is:\n"
        "- Concise and covers all the main points.\n"
        "- Avoids any hallucinations or fabricated information. Use only the provided details.\n"
        "- Accurate and factual, maintaining integrity throughout the summary.\n"
        "- Free of inappropriate language.\n"
        "- Summarize the document content, adding appropriate subheadings where necessary to enhance clarity and organization.\n"
        "- In the requested tone and format.\n"
        "- Provide a conclusion at the end."
    )

    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.3-70b-versatile",
    )

    return chat_completion.choices[0].message.content


# Function to generate content based on provided parameters
def generate_content(company_info, content_purpose, desired_action, topic_details, keywords, audience_profile, format_structure, num_words, seo_keywords, references):
    inputs = {
        "company_info": company_info,
        "content_purpose": content_purpose,
        "desired_action": desired_action,
        "topic_details": topic_details,
        "keywords": keywords,
        "audience_profile": audience_profile,
        "format_structure": format_structure,
        "seo_keywords": seo_keywords,
        "references": references
    }

    # Check if any input parameter contains inappropriate words
    inappropriate_key = None
    inappropriate_value = None
    for key, value in inputs.items():
        if value and contains_inappropriate_language(value):
            inappropriate_key = key
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in {inappropriate_key}: '{inappropriate_value}'."

    # Sanitize input fields
    sanitized_inputs = {key: sanitize_input(value) if value else '' for key, value in inputs.items()}

    # Construct the prompt for content generation
    prompt = f"Generate high-quality, engaging content of maximum {num_words} words with the following details:\n"

    if sanitized_inputs['company_info']:
        prompt += f"Company Information: {sanitized_inputs['company_info']}\n"
    if sanitized_inputs['content_purpose']:
        prompt += f"Purpose of Content: {sanitized_inputs['content_purpose']}\n"
    if sanitized_inputs['desired_action']:
        prompt += f"Desired Action: {sanitized_inputs['desired_action']}\n"
    if sanitized_inputs['topic_details']:
        prompt += f"Topic Details: {sanitized_inputs['topic_details']}\n"
    if sanitized_inputs['keywords']:
        prompt += f"Keywords: {sanitized_inputs['keywords']}\n"
    if sanitized_inputs['audience_profile']:
        prompt += f"Audience Profile: {sanitized_inputs['audience_profile']}\n"
    if sanitized_inputs['format_structure']:
        prompt += f"Format and Structure: {sanitized_inputs['format_structure']}\n"
    if sanitized_inputs['seo_keywords']:
        prompt += f"SEO Keywords: {sanitized_inputs['seo_keywords']}\n"
    if sanitized_inputs['references']:
        prompt += f"References to Cite: {sanitized_inputs['references']}\n"

    # Additional instructions for the content creation
    prompt += (
        "\nInstructions:\n"
        "1. Write a Head-Turning Headline: Create a headline that sparks interest, stirs emotion, or draws readers into the topic.\n"
        "2. Craft a Hook That Grabs Attention: Ensure the first sentence captures interest and transitions into the main point seamlessly.\n"
        "3. Incorporate Research and Data: Use credible statistics, metrics, or data to build authority and support your claims.\n"
        "4. Focus on a Single Purpose: Convey a clear, consistent key message throughout the content.\n"
        "5. Use a Unique Brand Voice: Align the tone and style of the content with the brand persona and target audience.\n"
        "6. Optimize for Digital: Format content with short paragraphs, bullet points, and use SEO best practices for better readability and search engine visibility.\n"
        "7. Edit for Perfection: Ensure the final content is polished, accurate, and free from errors through thorough editing.\n"
        "8. Engage and Inform: Highlight benefits and unique aspects to keep readers interested and informed.\n"
    )

    # Generate content using Groq API
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.3-70b-versatile",

    )

    return chat_completion.choices[0].message.content

def generate_sales_script(company_details, num_words, product_descriptions, features_benefits, pricing_info, promotions, target_audience, sales_objectives,
                          competitive_advantage, compliance):
    inputs = {
        "Company Details": company_details,
        "Product Descriptions": product_descriptions,
        "Features and Benefits": features_benefits,
        "Pricing Info": pricing_info,
        "Promotions": promotions,
        "Target Audience": target_audience,
        "Sales Objectives": sales_objectives,
        "Competitive Advantage": competitive_advantage,
        "Compliance": compliance,
        "Number Of Words": num_words
    }

    # Check if any input parameter contains inappropriate words
    inappropriate_key = None
    inappropriate_value = None
    for key, value in inputs.items():
        if value and contains_inappropriate_language(value):
            inappropriate_key = key
            inappropriate_value = value
            break

    if inappropriate_key:
        return f"This type of language is not allowed in {inappropriate_key}: {inappropriate_value}"

    # Sanitize input fields
    sanitized_inputs = {key: sanitize_input(value) if value else '' for key, value in inputs.items()}

    # Build the prompt for generating the sales script
    prompt = f"Generate a persuasive sales script of up to {sanitized_inputs['Number Of Words']} words using the following details:"
    for key, value in sanitized_inputs.items():
        if key != 'Number Of Words' and value:
            prompt += f"\n- {key}: {value}"

    prompt += (
        "\n\nInstructions:\n"
        "- Begin with a warm and professional introduction. Introduce yourself and the company, and ask politely if you can speak with the prospect.\n"
        "- Clearly explain what the company offers and how it aligns with the prospect's needs.\n"
        "- Include pre-qualifying questions that demonstrate curiosity about the prospect’s pain points or goals while keeping the tone conversational.\n"
        "- Proactively address common concerns or objections using empathetic and reassuring language to build trust.\n"
        "- Provide concise and compelling product information, focusing on features, benefits, and how they solve the prospect's challenges.\n"
        "- End with gratitude, a motivational call-to-action, and a follow-up suggestion to maintain engagement.\n"
        "- Ensure the tone remains empathetic, professional, and tailored to resonate with the prospect.\n"
        "- Maintain factual accuracy throughout the script and avoid any inappropriate language or fabricated details."
    )


    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)

    chat_completion = client.chat.completions.create(
        messages=[
            {
                "role": "user",
                "content": prompt,
            }
        ],
        model="llama-3.3-70b-versatile",

    )

    return chat_completion.choices[0].message.content


def bhashini_translate(text: str,  to_code: str = "Hindi", from_code: str = "English",user_id: str=BHASHINI_USER_ID, api_key: str=BHASHINI_API_KEY) -> dict:
    """Translates text from source language to target language using the Bhashini API.

    Args:
        text (str): The text to translate.
        from_code (str): Source language code. Default is 'en' (English).
        to_code (str): Target language code. Default is 'te' (Telugu).
        user_id (str): User ID for the API.
        api_key (str): API key for authentication.

    Returns:
        dict: A dictionary with the status code, message, and translated text or error info.
    """
    lang_dict = {
        "English": "en",
        "Hindi": "hi",
        "Tamil": "ta",
        "Telugu": "te",
        "Marathi": "mr",
        "Kannada": "kn",
        "Bengali": "bn",
        "Odia": "or",
        "Assamese": "as",
        "Punjabi": "pa",
        "Malayalam": "ml",
        "Gujarati": "gu",
        "Urdu": "ur",
        "Sanskrit": "sa",
        "Nepali": "ne",
        "Bodo": "brx",
        "Maithili": "mai",
        "Sindhi": "sd",
        "Kashmiri": "ks", 
        "Konkani": "kok",  
        "Dogri" :"doi",
        "Goan Konkani": "gom",
        "Santali": "sat"


    }


    from_code = lang_dict[from_code]
    to_code = lang_dict[to_code]



    # Setup the initial request to get model configurations
    url = 'https://meity-auth.ulcacontrib.org/ulca/apis/v0/model/getModelsPipeline'
    headers = {
        "Content-Type": "application/json",
        "userID": user_id,
        "ulcaApiKey": api_key
    }
    payload = {
        "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}}}],
        "pipelineRequestConfig": {"pipelineId": "64392f96daac500b55c543cd"}
    }
    response = requests.post(url, json=payload, headers=headers)

    if response.status_code != 200:
        return {"status_code": response.status_code, "message": "Error in translation request", "translated_content": None}

    # Process the response to setup the translation execution
    response_data = response.json()
    service_id = response_data["pipelineResponseConfig"][0]["config"][0]["serviceId"]
    callback_url = response_data["pipelineInferenceAPIEndPoint"]["callbackUrl"]
    headers2 = {
        "Content-Type": "application/json",
        response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["name"]: response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["value"]
    }
    compute_payload = {
        "pipelineTasks": [{"taskType": "translation", "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}, "serviceId": service_id}}],
        "inputData": {"input": [{"source": text}], "audio": [{"audioContent": None}]}
    }

    # Execute the translation
    compute_response = requests.post(callback_url, json=compute_payload, headers=headers2)
    if compute_response.status_code != 200:
        return {"status_code": compute_response.status_code, "message": "Error in translation", "translated_content": None}

    compute_response_data = compute_response.json()
    translated_content = compute_response_data["pipelineResponse"][0]["output"][0]["target"]

    return {"status_code": 200, "message": "Translation successful", "translated_content": translated_content}


def get_templates():
    return { 
        "default": "media/templates/default_theme_prod.pptx",
        "Professional Template 2": "media/templates/universal_blue_theme.pptx",
        "Professional Template 3": "media/templates/Liquid_void_presentation.pptx",
        "Professional Template 4": "media/templates/professional_business_theme.pptx",
        "Creative Template 1": "media/templates/artistic_fashion_final.pptx",
        "Creative Template 2": "media/templates/retro_theme.pptx",
        "Creative Template 3": "media/templates/japanese_floral.pptx",
        "Creative Template 4": "media/templates/peach_professional.pptx",
        "Futuristic Template 1": "media/templates/geometric_pink_theme.pptx",
        "Futuristic Template 2": "media/templates/blue_theme_cinematic.pptx",
        "Futuristic Template 3": "media/templates/futuristic_444.pptx",
        "Futuristic Template 4": "media/templates/3D_float_design_template.pptx",
        "Minimilistic Template 1": "media/templates/minimalist_lilac_theme.pptx",
        "Minimilistic Template 2": "media/templates/blue_spheres_theme.pptx",
        "Minimilistic Template 3": "media/templates/pastel_theme_futuristic.pptx",
        "Minimilistic Template 4": "media/templates/simple_yellow_theme.pptx",

    }

def generate_presentation_content(document_content, num_slides, special_instructions, title):
    """
    Generate all slide content in JSON format in a single API call.
    """
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
 
    # Prepare the prompt
    if document_content:
        prompt = (
            f"Based on the following document content  generate {num_slides} slide titles and their corresponding content:\n\n"
            f"{document_content}\n\n"
        )
    else:
        prompt = f"Based on the title '{title}', generate {num_slides} slide titles and their corresponding content.\n"
 
 
    if special_instructions:
        prompt += f" Pay attention to the following points: {special_instructions}. "
 
    prompt += (
        "Generate a JSON array of slide objects. "
        "Each object should have the keys 'title' and 'content'. "
        f"Constraint: Exactly {num_slides} slides objects should be generated always. "
        "The 'title' value must be a concise string, suitable for a slide title. "
        "The 'content' value must be a JSON array containing exactly four strings. "
        "Each string in the 'content' array must be a detailed point, between 15 and 20 words, suitable for a slide bullet point. "
        "Ensure that the content is professional, detailed, and not repetitive or gibberish. "
        "Avoid generating nested JSON structures. "
        "Do not include any additional text or explanations before and outside the JSON structure."
        "Do not include any additional text or explanations before and outside the array structure."

    )
    print("This is the prompt")
    # Fetch response from Groq API
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile", # Use specdec for better json formatting
    )
    print("This is the chat completion",chat_completion)
    generated_content = chat_completion.choices[0].message.content
    print("This is the generated content",generated_content)
 
 
    # Parse the JSON response
    try:
        slides_data = json.loads(generated_content)
        if not isinstance(slides_data, list):
            raise ValueError("Output is not a list of dictionaries")
        print(slides_data)
        return slides_data
    except json.JSONDecodeError as e:
         print(f"Failed to parse Groq response as JSON: {e}")
         print(f"Generated content: {generated_content}")
         raise ValueError(f"Failed to parse Groq response as JSON: {e}")

def update_presentation_with_generated_content(template_path, output_path,document_content,title,num_slides,special_instructions):
    """
    Update the PowerPoint presentation with the generated titles and content.
    """
    presentation = Presentation(template_path)
 
    slides_data = generate_presentation_content(
    document_content,num_slides,special_instructions,title)
 
    # Update slides with titles and content
    for slide_index, slide_data in enumerate(slides_data):
        if slide_index < len(presentation.slides) - 1:  # Exclude the last slide (e.g., "Thank You" slide)
            slide = presentation.slides[slide_index]
            i = 0
            for shape in slide.shapes:
                if shape.has_text_frame:
                    if i == 0:  # Title area
                        shape.text_frame.clear()
                        if slide_index == 0:  # Use user-entered title for first slide
                            shape.text_frame.text = title
                        else:  # Use AI-generated title for the rest
                            shape.text_frame.text = slide_data['title']
                    elif i == 1:  # Content area
                        shape.text_frame.clear()
                        shape.text_frame.text = '\n'.join(slide_data['content'])
                    i += 1
 
    # Remove extra slides, but preserve the last slide ("Thank You")
    while len(presentation.slides) > len(slides_data) + 1:
        slide_id = presentation.slides._sldIdLst[-2]  # Get the second-to-last slide ID
        presentation.slides._sldIdLst.remove(slide_id)  # Remove it from the slide list
 
    # Save the updated presentation
    presentation.save(output_path)
    print(f"Presentation saved as '{output_path}'")
    return presentation

def fetch_single_image(query, width, height): 
    headers = {"Authorization": PEXELS_API_KEY}
    params = {"query": query, "per_page": 10}  # Request up to 5 images
    response = requests.get(PEXELS_BASE_URL, headers=headers, params=params)
    if response.status_code == 200:
        data = response.json()
        photos = data.get("photos", [])
        if photos:
            random_photo = random.choice(photos)  # Ensure 'photos' is a list
            original_url = random_photo['src']['landscape']
            resized_url = f"{original_url}?auto=compress&cs=tinysrgb&h={height}&w={width}"

            # Fetch the image data
            image_response = requests.get(resized_url)
            if image_response.status_code == 200:
                image_data = base64.b64encode(image_response.content).decode('utf-8')  # Encode image as base64
                return {
                    "url": resized_url,
                    "base64_image": image_data
                }
            else:
                print(f"Error fetching image: {image_response.status_code}")
                return None
        else:
            print("No images found for the given query.")
            return None
    else:
        print(f"Error: {response.status_code} - {response.text}")
        return None
    
def generate_blog(title, tone, custom_tone, keywords=None):
    # Ensure all fields are checked for inappropriate language
    fields_to_check = [title, keywords]
    if any(contains_inappropriate_language(str(field)) for field in fields_to_check if field is not None):
        return "Error: Input contains inappropriate language."
    print(keywords)
    

    # Create the prompt for the blog generation
    # prompt = (
    #     f"As a professional blog generator with extensive experience in the field, create a compelling blog post "
    #     f"with a maximum of {1000} words. The title of the blog is '{title}' and it should be written in a "
    #     f"{tone} tone with a custom tone of '{custom_tone}' and appropriate side headings."
    # )
    # prompt += (
    #     "\n\nMake sure to introduce the topic engagingly and provide valuable insights throughout the blog. "
    #     "Include relevant examples or case studies to illustrate key points."
    # )

    # if keywords:
    #     prompt += (
    #         f"\n\nIncorporate the following keywords: {', '.join(keywords)}. Ensure that the keywords are naturally "
    #         "integrated into the content to enhance readability and SEO."
    #     )

    # prompt += (
    #     "\n\nConclude the blog with a strong closing statement that summarizes the key points."
    #   #  "\n\nAlso, include a call-to-action or thought-provoking idea to engage the reader."
    #     "\n\nEnsure that the content is informative, easy to understand, and engaging. Avoid jargon and ensure that "
    #     "the blog flows logically from one section to the next."
    # )
     


    prompt = (
        f"I want you to craft a professional, engaging, and commercially valuable blog. Use the following details provided by the user to craft a well-structured and impactful blog:\n\n"
        f"1. **Title**: {title}\n"
        f"2. **Keywords**: {', '.join(keywords) if keywords else 'N/A'}\n"
        f"3. **Tone**: {tone} (custom tone: {custom_tone})\n\n"
        f"**Instructions:**\n"
        f"- Start with an attention-grabbing introduction that hooks the reader and clearly outlines the topic and its relevance.\n"
        f"- Organize the content into clear sections with subheadings for readability.\n"
        f"- Incorporate the provided keywords naturally throughout the blog to enhance SEO without making it feel forced.\n"
        f"- Ensure each section provides valuable insights, actionable advice, or unique perspectives that align with the title.\n"
        f"- Use persuasive language, real-world examples, and credible data where appropriate to enhance the content's appeal and trustworthiness.\n"
        f"- Conclude with a thought-provoking summary or call to action that encourages further engagement (e.g., sharing the blog, visiting a website, or purchasing a product).\n"
        f"- Maintain a tone that matches the provided tone style, keeping the target audience in mind.\n\n"
        f"The blog should be comprehensive, well-researched, and strictly ensuring in a range of 900-1100 words  to maximize its impact and commercial value."

    )
    # Call the Groq API to generate the blog
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile",
    )

    # Return the generated blog content
    return chat_completion.choices[0].message.content

def rephrasely(text_to_rephrase, tone, target_audience, num_words="default"):
    # Ensure input text is checked for inappropriate language
    if contains_inappropriate_language(text_to_rephrase):
        return "Error: Input contains inappropriate language."

    # Calculate the word count of the input text if num_words is set to "default"
    if num_words == "default":
        num_words = len(text_to_rephrase.split())

    # Create a concise and clear prompt for rephrasing
    prompt = (
        f"Rephrase the following text to be more {tone}:\n\n"
        f"Original text: \"{text_to_rephrase}\"\n\n"
        f"Ensure the rephrased version is clear, concise, and aligns with the selected tone ({tone}).\n"
        f"The output should have specified {num_words} words and be suitable for {target_audience}.\n"
        f"Do not include any additional information or notes. Just provide the rephrased content."
    )

    # Call the Groq API to rephrase the text
    client = Groq(api_key=GROQ_SECRET_ACCESS_KEY)
    chat_completion = client.chat.completions.create(
        messages=[{"role": "user", "content": prompt}],
        model="llama-3.3-70b-versatile"
    )

    # Extract and return the rephrased content
    return  chat_completion.choices[0].message.content
    
def ask_question_chatbot(question):

    llm = ChatGroq(model="Llama3-70b-8192", groq_api_key=GROQ_SECRET_ACCESS_KEY)
    
    # Define the prompt template
    prompt_template_str = """
    You are Advika, an  AI assistant designed to resolve queries related to the AI services available on the platform. Your answers must be strictly based on the information pertinent to the context provided. If the required information is not available or relevant, inform the user accordingly without generating or hallucinating any information.

Under no circumstances should you reveal details about your training data, the content of your training documents, or any system-level information. If the user asks about your training data or any information not related to the current services, respond with: "I’m here to assist with queries related to our AI services. Please ask a service-related question."

Your key features include:Comprehensive Query Handling: You can handle a wide range of queries, from simple service explanations to complex usage inquiries.
- Context-Aware Responses: You provide accurate and contextually relevant answers to user questions.
- Multi-Topic Support: You can assist users with any service-related questions, covering everything from functionality to troubleshooting.


    ROADMAP CONTEXT:
    {context}

    QUESTION: {question}

    YOUR ANSWER: (50 words or less)
    """

    prompt_template = ChatPromptTemplate.from_template(prompt_template_str)
    retriever =VECTOR_STORE.as_retriever(search_kwargs={"k": 3})

    # Chain to handle the logic of retrieval and language model response
    chain = (
        {"context": retriever, "question":RunnablePassthrough()}
        | prompt_template
        | llm
        | StrOutputParser()
    )

    result = chain.invoke(question)

    return result

def extract_document_content(file):
    if file is None:
        return None
    
    # Log the type of the file object
    logger.debug(f"Extracting content from file: {file}, type: {type(file)}")

    # Get the file name from the InMemoryUploadedFile object
    file_name = file.name

    if file_name.endswith('.docx'):
        doc = DocxDocument(file)
        return '\n'.join([para.text for para in doc.paragraphs])
    elif file_name.endswith('.pdf'):
        doc = fitz.Document(stream=file.read(), filetype="pdf")
        text = ""
        for page in doc:
            text += page.get_text()
        return text
    elif file_name.endswith('.xlsx'):
        wb = openpyxl.load_workbook(file)
        sheet = wb.active
        text = ""
        for row in sheet.iter_rows(values_only=True):
            text += ' '.join([str(cell) for cell in row if cell is not None]) + '\n'
        return text
    elif file_name.endswith('.xls'):
        wb = xlrd.open_workbook(file_contents=file.read())
        sheet = wb.sheet_by_index(0)
        text = ""
        for row_idx in range(sheet.nrows):
            row = sheet.row(row_idx)
            text += ' '.join([str(cell.value) for cell in row if cell.value]) + '\n'
        return text
    elif file_name.endswith('.pptx'):
        # Extract content from PPTX file using python-pptx
        text = ""
        presentation = Presentation(file)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    elif file_name.endswith('.ppt'):
        # Extract content from PPT file (simplest way is to convert to PPTX or use another library)
        # If you can't handle .ppt directly, you might want to raise an error or return an informative message.
        raise ValueError("Unsupported file type: .ppt files need to be converted to .pptx.")
    else:
        raise ValueError("Unsupported file type")

import requests
from typing import List, Dict
from concurrent.futures import ThreadPoolExecutor, as_completed
import time
# Cache for storing translation results
translation_cache = {}
def retry(func, retries=3, delay=1):
    """Retries a function call a specified number of times with a delay."""
    for attempt in range(retries):
        try:
            return func()
        except Exception as e:
            if attempt < retries - 1:
                time.sleep(delay)
                continue
            else:
                raise e
def bhashini_translate_json(text: str, to_code: str = "Hindi", from_code: str = "English",
                       user_id: str = BHASHINI_USER_ID, api_key: str = BHASHINI_API_KEY) -> dict:
    """Translates text from source language to target language using the Bhashini API."""
    
    # Cache key based on input parameters
    cache_key = f"{text}_{from_code}_{to_code}"
    if cache_key in translation_cache:
        return translation_cache[cache_key]
    lang_dict = {
        "English": "en", "Hindi": "hi", "Tamil": "ta", "Telugu": "te",
        "Marathi": "mr", "Kannada": "kn", "Bengali": "bn", "Odia": "or",
        "Assamese": "as", "Punjabi": "pa", "Malayalam": "ml", "Gujarati": "gu",
        "Urdu": "ur", "Sanskrit": "sa", "Nepali": "ne", "Bodo": "brx",
        "Maithili": "mai", "Sindhi": "sd", "Kashmiri": "ks", "Konkani": "kok",
        "Dogri": "doi", "Goan Konkani": "gom", "Santali": "sat"
    }
    from_code = lang_dict[from_code]
    to_code = lang_dict[to_code]
    # Setup the initial request to get model configurations
    url = 'https://meity-auth.ulcacontrib.org/ulca/apis/v0/model/getModelsPipeline'
    headers = {
        "Content-Type": "application/json",
        "userID": user_id,
        "ulcaApiKey": api_key
    }
    payload = {
        "pipelineTasks": [{
            "taskType": "translation",
            "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}}
        }],
        "pipelineRequestConfig": {"pipelineId": "64392f96daac500b55c543cd"}
    }
    
    # Retry the model configuration request
    response = retry(lambda: requests.post(url, json=payload, headers=headers))
    if response.status_code != 200:
        return {"status_code": response.status_code, "message": "Error in translation request", "translated_content": None}
    # Process the response to setup the translation execution
    response_data = response.json()
    service_id = response_data["pipelineResponseConfig"][0]["config"][0]["serviceId"]
    callback_url = response_data["pipelineInferenceAPIEndPoint"]["callbackUrl"]
    headers2 = {
        "Content-Type": "application/json",
        response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["name"]: response_data["pipelineInferenceAPIEndPoint"]["inferenceApiKey"]["value"]
    }
    compute_payload = {
        "pipelineTasks": [{
            "taskType": "translation",
            "config": {"language": {"sourceLanguage": from_code, "targetLanguage": to_code}, "serviceId": service_id}
        }],
        "inputData": {"input": [{"source": text}], "audio": [{"audioContent": None}]}
    }
    # Retry the translation execution
    compute_response = retry(lambda: requests.post(callback_url, json=compute_payload, headers=headers2))
    if compute_response.status_code != 200:
        return {"status_code": compute_response.status_code, "message": "Error in translation", "translated_content": None}
    compute_response_data = compute_response.json()
    translated_content = compute_response_data["pipelineResponse"][0]["output"][0]["target"]
    # Cache the translation result
    translation_cache[cache_key] = {"status_code": 200, "message": "Translation successful", "translated_content": translated_content}
    
    return translation_cache[cache_key]
def translate_multiple_texts(text_list: List[str], from_code: str, to_code: str,
                              user_id: str, api_key: str) -> List[dict]:
    """Translates multiple texts concurrently using Bhashini API."""
    results = []
    with ThreadPoolExecutor(max_workers=100) as executor:
        future_to_text = {executor.submit(bhashini_translate_json, text, to_code, from_code, user_id, api_key): text for text in text_list}
        for future in as_completed(future_to_text):
            text = future_to_text[future]
            try:
                result = future.result()
                results.append(result)
            except Exception as e:
                print(f"Translation failed for {text}: {e}")
                results.append({"text": text, "error": str(e)})
    return results
